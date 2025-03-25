from docx import Document
from openpyxl import Workbook
from pptx import Presentation

# Word Document
doc = Document()
doc.add_heading('Sample Document', 0)
doc.add_paragraph('This is a test document created with python-docx.')
doc.save('sample.docx')

# Excel Spreadsheet
wb = Workbook()
ws = wb.active
ws['A1'] = 'Hello'
ws['B1'] = 'World'
wb.save('sample.xlsx')

# PowerPoint Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = 'Sample Slide'
prs.save('sample.pptx')

print("Files created: sample.docx, sample.xlsx, sample.pptx")