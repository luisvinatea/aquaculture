from pptx import Presentation

# Paths to the PPTX and the output HTML file
pptx_path = "/mnt/raid_cache/Data/Gdrive/aquaculture/conaqua/docs/presentations/conacua.pptx"
output_html_path = "conacua.html"

# Load the PPTX
prs = Presentation(pptx_path)

# Create an HTML representation of the presentation
with open(output_html_path, "w", encoding="utf-8") as html_file:
    html_file.write("<html><body>\n")
    
    for slide_idx, slide in enumerate(prs.slides):
        html_file.write(f"<div class='slide' id='slide-{slide_idx}'>\n")
        html_file.write(f"<h2>Slide {slide_idx + 1}</h2>\n")
        
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_frame = shape.text_frame
                
                # Write each paragraph within the shape as a separate element
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    cleaned_text = paragraph.text.strip()
                    if cleaned_text:
                        html_file.write(f"<p id='slide-{slide_idx}-shape-{shape_idx}-para-{para_idx}'>{cleaned_text}</p>\n")
        
        html_file.write("</div>\n")
    
    html_file.write("</body></html>")

print(f"HTML representation of the PPTX saved to {output_html_path}")