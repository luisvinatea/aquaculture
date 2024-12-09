from bs4 import BeautifulSoup
from pptx import Presentation

# Paths to the PPTX, edited HTML file, and the output PPTX
pptx_path = "/mnt/raid_cache/Data/Gdrive/aquaculture/conaqua/docs/presentations/conacua.pptx"
html_path = "conacua.html"
output_pptx_path = "/mnt/raid_cache/Data/Gdrive/aquaculture/conaqua/docs/presentations/conacua_translated.pptx"

# Load the HTML with translations
with open(html_path, "r", encoding="utf-8") as html_file:
    soup = BeautifulSoup(html_file, "html.parser")

# Load the original PPTX
prs = Presentation(pptx_path)

# Function to update PPTX content with translations from the HTML
for slide_idx, slide in enumerate(prs.slides):
    slide_div = soup.find("div", {"id": f"slide-{slide_idx}"})
    
    if slide_div:
        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            paragraphs = text_frame.paragraphs
            
            for para_idx, paragraph in enumerate(paragraphs):
                para_tag = slide_div.find("p", {"id": f"slide-{slide_idx}-shape-{shape_idx}-para-{para_idx}"})
                
                if para_tag:
                    paragraph.text = para_tag.get_text(strip=True)

# Save the updated PPTX
prs.save(output_pptx_path)
print(f"Translated PPTX saved at {output_pptx_path}")